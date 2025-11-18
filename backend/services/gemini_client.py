
import json
import os
from typing import List, Dict, Any

from google import genai
# Add pptx support
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")


class GeminiClient:
    def __init__(self):
        self.client = genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None
        print("GEMINI_API_KEY:", GEMINI_API_KEY)
        # You can change model names if needed
        self.text_model_name = "gemini-2.5-flash"
        self.embed_model_name = "text-embedding-004"

    def _ensure_client(self) -> bool:
        return self.client is not None

    def extract_text_from_file(self, file_obj: Dict[str, Any]) -> str:
        """
        Use Gemini to turn an uploaded file (PDF, doc, image, etc.)
        into plain text for RAG. If file is .pptx, extract text using python-pptx first, then send to Gemini.
        """
        if not self._ensure_client():
            return ""

        mime_type = file_obj.get("content_type") or "application/octet-stream"
        data = file_obj["data"]

        # If PowerPoint, extract text using python-pptx, then send to Gemini
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            if Presentation is None:
                return "Error: python-pptx is not installed. Cannot process .pptx files."
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=True) as tmp:
                tmp.write(data)
                tmp.flush()
                prs = Presentation(tmp.name)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                extracted_text = "\n".join(text)
            # Now send extracted text to Gemini
            prompt = (
                "You are a document text extractor. "
                "Read the following extracted text from a PowerPoint file and return ONLY the plain text content, "
                "with no formatting, explanations, or extra commentary."
            )
            contents = [
                {
                    "role": "user",
                    "parts": [
                        {"text": extracted_text},
                        {"text": prompt},
                    ],
                }
            ]
            response = self.client.models.generate_content(
                model=self.text_model_name,
                contents=contents,
            )
            return (getattr(response, "text", "") or "").strip()

        # Otherwise, use Gemini for supported types
        prompt = (
            "You are a document text extractor. "
            "Read the content of the attached file and return ONLY the plain text content, "
            "with no formatting, explanations, or extra commentary."
            "you can recive files can be image, pdf, docx, txt, video etc."
        )

        contents = [
            {
                "role": "user",
                "parts": [
                    {
                        "inline_data": {
                            "mime_type": mime_type,
                            "data": data,
                        }
                    },
                    {"text": prompt},
                ],
            }
        ]

        response = self.client.models.generate_content(
            model=self.text_model_name,
            contents=contents,
        )

        return (getattr(response, "text", "") or "").strip()

    def chunk_text_for_rag(self, text: str, max_chars: int = 800) -> List[Dict[str, str]]:
        chunks: List[Dict[str, str]] = []
        current: List[str] = []
        current_len = 0

        for paragraph in text.split("\n"):
            p = paragraph.strip()
            if not p:
                continue
            if current_len + len(p) + 1 > max_chars and current:
                chunks.append({"text": "\n".join(current)})
                current = [p]
                current_len = len(p)
            else:
                current.append(p)
                current_len += len(p) + 1

        if current:
            chunks.append({"text": "\n".join(current)})

        return chunks

    def embed_text(self, text: str) -> List[float]:
        if not self._ensure_client():
            return []

        resp = self.client.models.embed_content(
            model=self.embed_model_name,
            contents=[text],
        )

        # Adapt to common response shape: embeddings[0].values
        try:
            return list(resp.embeddings[0].values)
        except Exception:
            return []

    def answer_with_context(
        self,
        question: str,
        context_chunks: List[Dict[str, Any]],
        review_threshold: float = 0.6,
    ) -> Dict[str, Any]:
        """
        Use Gemini to answer a question using ONLY the provided context chunks.
        Returns: { answer, confidence, citations, needs_human_review }
        """
        if not self._ensure_client():
            return {
                "answer": "",
                "confidence": 0.0,
                "citations": [],
                "needs_human_review": True,
                "error": "GEMINI_API_KEY not configured",
            }

        if not context_chunks:
            return {
                "answer": "I cannot answer this because there is no relevant context yet. Please upload some documents first.",
                "confidence": 0.0,
                "citations": [],
                "needs_human_review": True,
            }

        context_lines: List[str] = []
        for i, ch in enumerate(context_chunks):
            # ContextChunk may be a Pydantic model; normalize to dict-like access.
            if hasattr(ch, "dict"):
                ch_data = ch.dict()
            elif isinstance(ch, dict):
                ch_data = ch
            else:
                ch_data = {
                    "source": getattr(ch, "source", ""),
                    "chunk_index": getattr(ch, "chunk_index", -1),
                    "text": getattr(ch, "text", ""),
                }

            context_lines.append(
                f"[{i}] source={ch_data.get('source','')}, chunk_index={ch_data.get('chunk_index', -1)}\n{ch_data.get('text','')}\n"
            )
        context_text = "\n\n".join(context_lines)

        system_prompt = """
You are an answer-generation agent for a Retrieval-Augmented Generation (RAG) system.

You are given:
- a user's question
- a set of context chunks extracted from the user's private documents

Rules:
1. Use ONLY the provided context to answer. Do NOT use outside knowledge.
2. If the answer is not clearly supported by the context, say:
   "I cannot answer this based on the provided context."
3. Be concise and factual.
4. Provide a confidence score between 0.0 and 1.0 that reflects how well the context supports your answer, so if you are not sure with you answer or the question than just give me actual confidence score.
5. Provide citations: a list of objects { "source": string, "chunk_index": number } corresponding to the chunks you used.
6. if the confidence score is below 0.6, needs_human_review must be true.

Return ONLY valid JSON in this exact format:

{
  "answer": "<string>",
  "confidence": <float between 0 and 1>,
  "citations": [
    { "source": "<string>", "chunk_index": <number> }
  ],
  "needs_human_review": <boolean>
}
"""

        prompt_text = (
            system_prompt
            + "\n\nQuestion:\n"
            + question
            + "\n\nContext Chunks:\n"
            + context_text
        )

        response = self.client.models.generate_content(
            model=self.text_model_name,
            contents=prompt_text,
        )

        raw = (getattr(response, "text", "") or "").strip()

        # Strip code block wrapper if present (```json ... ```)
        if raw.startswith("```"):
            lines = raw.split("\n")
            # Remove first line (```json or similar) and last line (```)
            if len(lines) > 2:
                raw = "\n".join(lines[1:-1])
            elif lines[0].startswith("```"):
                raw = lines[0][3:]  # Remove leading ```
            if raw.startswith("json"):
                raw = raw[4:].strip()

        try:
            parsed = json.loads(raw)
            answer = parsed.get("answer", "")
            try:
                confidence = float(parsed.get("confidence", 0.0) or 0.0)
            except Exception:
                confidence = 0.0
            citations = parsed.get("citations", [])
        except json.JSONDecodeError:
            # Fallback: assign confidence=0.5
            answer = raw
            confidence = 0.5
            citations = []

        # Always check if review is needed based on threshold
        needs_human_review = confidence < review_threshold

        result = {
            "answer": answer,
            "confidence": confidence,
            "citations": citations,
            "needs_human_review": needs_human_review,
        }

        # If human review is needed, generate a context-aware follow-up question using LLM
        if needs_human_review:
            followup_question = self.generate_followup_question(question, context_chunks, answer)
            result["followup_question"] = followup_question

        return result

    def generate_followup_question(self, question: str, context_chunks: List[Dict[str, Any]], answer: str = "") -> str:
        """
        Use Gemini LLM to generate a context-aware follow-up question for the user if human review is needed.
        """
        if not self._ensure_client():
            return "Can you clarify or provide more details?"

        # Build context for LLM
        context_lines = []
        for i, ch in enumerate(context_chunks):
            context_lines.append(
                f"[{i}] source={ch.get('source','')}, chunk_index={ch.get('chunk_index', -1)}\n{ch.get('text','')}\n"
            )
        context_text = "\n\n".join(context_lines)

        prompt = (
            "You are an assistant helping a user with document Q&A. "
            "The previous answer is the model's best guess but may be incomplete or unreliable. "
            "Your job is to ask one natural-sounding follow-up question that helps clarify the user's original question so the next answer can be more accurate. "
            "Base the follow-up on the mismatch between the user's question and the previous answer, and use the context chunks only as needed for grounding. "
            "Ask exactly one question in plain language.\n"
            f"Original question: {question}\n"
            f"Previous answer: {answer}\n"
            f"Context chunks:\n{context_text}"
        )
        contents = [
            {
                "role": "user",
                "parts": [
                    {"text": prompt},
                ],
            }
        ]
        response = self.client.models.generate_content(
            model=self.text_model_name,
            contents=contents,
        )
        return (getattr(response, "text", "") or "Can you clarify your question?").strip()
